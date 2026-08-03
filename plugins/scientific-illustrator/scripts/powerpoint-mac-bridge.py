#!/usr/bin/env python3
"""Cross-platform file-backed PPTX bridge for Scientific Illustrator.

The preferred Windows Microsoft PowerPoint backend edits the live COM model.
This bridge covers Microsoft PowerPoint for Mac and WPS Presentation on Windows
or macOS. It builds native editable OOXML objects with python-pptx, keeps them
in an isolated working copy, and reopens that copy in the selected presentation
application after each mutation. Background refresh preserves the user's current
foreground application by default. The system mouse and keyboard are never used.
"""

from __future__ import annotations

import base64
import copy
import json
import math
import os
from pathlib import Path
import shutil
import subprocess
import sys
import tempfile
import time
import uuid

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt


STATE_DIR = Path(os.environ.get(
    "SCIENTIFIC_ILLUSTRATOR_MAC_DIR",
    str(Path.home() / ".codex" / "scientific-illustrator" / "macos"),
)).expanduser()
STATE_PATH = STATE_DIR / "session.json"
MAC_HOSTS = {
    "powerpoint": [Path("/Applications/Microsoft PowerPoint.app")],
    "wps": [
        Path("/Applications/wpsoffice.app"),
        Path("/Applications/WPS Office.app"),
        Path("/Applications/WPS.app"),
    ],
}


def _requested_host(args: dict | None = None, state: dict | None = None) -> str:
    explicit = (args or {}).get("host_application")
    configured = os.environ.get("SCIENTIFIC_ILLUSTRATOR_PPT_HOST", "auto")
    saved = (state or {}).get("host_application")
    value = str(explicit or saved or configured or "auto").strip().lower()
    if value not in {"auto", "powerpoint", "wps"}:
        raise ValueError("host_application must be auto, powerpoint, or wps")
    return value


def _windows_wps_candidates() -> list[Path]:
    candidates = []
    if os.environ.get("WPS_PRESENTATION_PATH"):
        candidates.append(Path(os.environ["WPS_PRESENTATION_PATH"]))
    for variable in ("ProgramFiles", "ProgramFiles(x86)", "LOCALAPPDATA"):
        root = os.environ.get(variable)
        if not root:
            continue
        base = Path(root)
        candidates.extend([
            base / "Kingsoft" / "WPS Office" / "office6" / "wpp.exe",
            base / "Kingsoft" / "WPS Office" / "office6" / "wpsoffice.exe",
        ])
        candidates.extend((base / "Kingsoft" / "WPS Office").glob("*/office6/wpp.exe"))
        candidates.extend((base / "Kingsoft" / "WPS Office").glob("*/office6/wpsoffice.exe"))
    return candidates


def _available_hosts() -> dict[str, dict]:
    if sys.platform == "darwin":
        result = {}
        for key, candidates in MAC_HOSTS.items():
            found = next((item for item in candidates if item.exists()), None)
            result[key] = {"installed": found is not None, "path": str(found) if found else None}
        return result
    if sys.platform == "win32":
        wps = next((item for item in _windows_wps_candidates() if item.exists()), None)
        powerpoint_installed = False
        powershell = shutil.which("powershell.exe") or shutil.which("powershell")
        if powershell:
            probe = _run([powershell, "-NoProfile", "-Command", "if (Test-Path 'Registry::HKEY_CLASSES_ROOT\\PowerPoint.Application\\CLSID') { 'yes' }"], check=False)
            powerpoint_installed = "yes" in probe.stdout.lower()
        return {
            "powerpoint": {"installed": powerpoint_installed, "path": None},
            "wps": {"installed": wps is not None, "path": str(wps) if wps else None},
        }
    return {"powerpoint": {"installed": False, "path": None}, "wps": {"installed": False, "path": None}}


def _select_host(args: dict | None = None, state: dict | None = None) -> tuple[str, dict]:
    requested = _requested_host(args, state)
    available = _available_hosts()
    if requested == "auto":
        for candidate in ("powerpoint", "wps"):
            if available[candidate]["installed"]:
                return candidate, available[candidate]
        return "powerpoint", available["powerpoint"]
    return requested, available[requested]


def _state() -> dict:
    if STATE_PATH.exists():
        try:
            return json.loads(STATE_PATH.read_text("utf-8"))
        except Exception:
            pass
    return {"path": None, "source_path": None, "metadata": {}, "owned": False}


def _write_state(state: dict) -> None:
    STATE_DIR.mkdir(parents=True, exist_ok=True)
    STATE_PATH.write_text(json.dumps(state, ensure_ascii=False, indent=2), "utf-8")


def _run(command: list[str], *, check: bool = True, timeout: int = 120) -> subprocess.CompletedProcess:
    return subprocess.run(command, check=check, capture_output=True, text=True, timeout=timeout)


def _osascript(source: str, *arguments: str, check: bool = False) -> str:
    result = _run(["osascript", "-e", source, *arguments], check=check, timeout=30)
    return result.stdout.strip()


def _focus_policy() -> str:
    value = os.environ.get("SCIENTIFIC_ILLUSTRATOR_FOCUS_POLICY", "preserve").strip().lower()
    return value if value in {"preserve", "foreground"} else "preserve"


def _open_windows_presentation(file_path: Path, executable: str | None, focus_policy: str) -> None:
    if focus_policy == "foreground":
        if executable:
            subprocess.Popen([executable, str(file_path)], close_fds=True)
        else:
            os.startfile(str(file_path))
        return

    # SW_SHOWNOACTIVATE prevents a newly created WPS/PowerPoint window from
    # taking focus. Some already-running hosts ignore that hint, so restore the
    # previously focused window after dispatching the file-open request.
    import ctypes

    user32 = ctypes.windll.user32
    previous_foreground = user32.GetForegroundWindow()
    if executable:
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = 4  # SW_SHOWNOACTIVATE
        subprocess.Popen([executable, str(file_path)], close_fds=True, startupinfo=startupinfo)
    else:
        result = ctypes.windll.shell32.ShellExecuteW(None, "open", str(file_path), None, None, 4)
        if result <= 32:
            raise RuntimeError(f"Unable to open presentation in the background; ShellExecuteW returned {result}.")
    if previous_foreground:
        time.sleep(0.15)
        user32.SetForegroundWindow(previous_foreground)


def _presentation_host_info(args: dict | None = None, state: dict | None = None) -> dict:
    host_name, host = _select_host(args, state)
    installed = bool(host["installed"])
    version = ""
    presentation_count = 0
    if installed and sys.platform == "darwin" and host_name == "powerpoint":
        version = _osascript('tell application "Microsoft PowerPoint" to get version')
        count_text = _osascript('tell application "Microsoft PowerPoint" to get count of presentations')
        try:
            presentation_count = int(count_text)
        except Exception:
            presentation_count = 0
    process_ids: list[int] = []
    if sys.platform == "darwin":
        process_pattern = "/Microsoft PowerPoint.app/" if host_name == "powerpoint" else "(wpsoffice|WPS Office|WPS).app"
        pgrep = _run(["pgrep", "-f", process_pattern], check=False)
        process_ids = [int(value) for value in pgrep.stdout.split() if value.isdigit()]
    elif sys.platform == "win32":
        image_names = ["POWERPNT.EXE"] if host_name == "powerpoint" else ["wpp.exe", "wpsoffice.exe"]
        for image_name in image_names:
            tasklist = _run(["tasklist", "/FI", f"IMAGENAME eq {image_name}", "/FO", "CSV", "/NH"], check=False)
            for line in tasklist.stdout.splitlines():
                fields = [item.strip('"') for item in line.split('","')]
                if len(fields) > 1 and fields[1].isdigit():
                    process_ids.append(int(fields[1]))
    return {
        "platform": sys.platform,
        "backend": "python-pptx-ooxml+application-reload",
        "host_application": host_name,
        "available_hosts": _available_hosts(),
        "application_path": host.get("path"),
        "installed": installed,
        "application_version": version,
        "running_processes": len(process_ids),
        "process_ids": process_ids,
        "active_application_process_id": process_ids[0] if process_ids else 0,
        "presentation_count": presentation_count,
    }


def _refresh_presentation(file_path: Path, state: dict | None = None, *, focus_policy: str | None = None) -> None:
    if os.environ.get("SCIENTIFIC_ILLUSTRATOR_POWERPOINT_SYNC", "1") == "0":
        return
    host_name, host = _select_host(state=state)
    if not host["installed"]:
        return
    focus_policy = focus_policy or _focus_policy()
    if sys.platform == "win32":
        _open_windows_presentation(file_path, host.get("path"), focus_policy)
        return
    if host_name == "powerpoint":
        close_script = r'''
on run argv
    set targetName to item 1 of argv
    tell application "Microsoft PowerPoint"
        repeat with deck in presentations
            try
                if (name of deck as text) is targetName then close deck saving no
            end try
        end repeat
    end tell
end run
'''
        try:
            _run(["osascript", "-e", close_script, file_path.name], check=False, timeout=8)
        except subprocess.TimeoutExpired:
            # A modal/read-only dialog must never make an MCP edit hang. The
            # managed file is still safely updated on disk and can be reopened.
            pass
    # `open -g -a` returns immediately, works for both Microsoft PowerPoint and
    # WPS, and preserves the user's foreground application. Foreground mode
    # omits `-g`. Direct AppleScript `open` can block indefinitely when
    # PowerPoint has a modal/read-only document, so it is not used here.
    open_args = ["open"]
    if focus_policy == "preserve":
        open_args.append("-g")
    open_args.extend(["-a", host["path"], str(file_path)])
    subprocess.Popen(open_args, close_fds=True)


def _managed_path(label: str = "scientific-illustrator") -> Path:
    STATE_DIR.mkdir(parents=True, exist_ok=True)
    return STATE_DIR / f"{label}-{time.strftime('%Y%m%d-%H%M%S')}-{uuid.uuid4().hex[:6]}.pptx"


def _new_presentation(path: Path) -> Presentation:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return prs


def _require_path(create: bool = False) -> tuple[dict, Path]:
    state = _state()
    raw_path = state.get("path")
    if raw_path:
        path = Path(raw_path)
        if path.exists():
            return state, path
    if not create:
        raise RuntimeError("No managed PowerPoint presentation is active. Call powerpoint_launch or powerpoint_new_presentation first.")
    path = _managed_path()
    _new_presentation(path)
    host_name, _ = _select_host(state=state)
    state.update({"path": str(path), "source_path": None, "owned": True, "metadata": {}, "host_application": host_name})
    _write_state(state)
    return state, path


def _load() -> tuple[dict, Path, Presentation]:
    state, path = _require_path(False)
    return state, path, Presentation(path)


def _save(prs: Presentation, state: dict, path: Path, *, refresh: bool = True) -> None:
    temp_path = path.with_suffix(".saving.pptx")
    prs.save(temp_path)
    os.replace(temp_path, path)
    _write_state(state)
    if refresh:
        _refresh_presentation(path, state)


def _slide(prs: Presentation, index: int):
    if index < 1 or index > len(prs.slides):
        raise ValueError(f"slide_index {index} is outside 1..{len(prs.slides)}")
    return prs.slides[index - 1]


def _shape(slide, args: dict):
    wanted_name = args.get("shape_name")
    wanted_id = args.get("shape_id")
    for shape in slide.shapes:
        if wanted_name is not None and shape.name.lower() == str(wanted_name).lower():
            return shape
        if wanted_id is not None and shape.shape_id == int(wanted_id):
            return shape
    raise ValueError(f"Shape not found: {wanted_name or wanted_id}")


def _pt(value: float | int) -> Pt:
    return Pt(float(value))


def _rgb(value: str | None) -> RGBColor | None:
    if not value:
        return None
    text = str(value).strip().lstrip("#")
    if len(text) != 6:
        raise ValueError(f"Invalid color {value}; expected #RRGGBB")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _set_fill(shape, color: str | None) -> None:
    rgb = _rgb(color)
    if rgb is None:
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_line(shape, args: dict) -> None:
    if not hasattr(shape, "line"):
        return
    rgb = _rgb(args.get("line_color"))
    if rgb is not None:
        shape.line.color.rgb = rgb
    if args.get("line_width") is not None:
        shape.line.width = _pt(args["line_width"])


def _set_arrow(shape, start_arrow: str = "none", end_arrow: str = "none") -> None:
    if not hasattr(shape, "line"):
        return
    line = shape.line._get_or_add_ln()
    for tag, kind in (("a:headEnd", start_arrow), ("a:tailEnd", end_arrow)):
        existing = line.find(tag, line.nsmap)
        if existing is None:
            from pptx.oxml.xmlchemy import OxmlElement
            existing = OxmlElement(tag)
            line.append(existing)
        mapping = {"none": "none", "open": "arrow", "triangle": "triangle", "stealth": "stealth", "diamond": "diamond", "oval": "oval"}
        existing.set("type", mapping.get(str(kind), "none"))


def _apply_text_frame(text_frame, args: dict) -> None:
    if args.get("margin_left") is not None:
        text_frame.margin_left = _pt(args["margin_left"])
    if args.get("margin_right") is not None:
        text_frame.margin_right = _pt(args["margin_right"])
    if args.get("margin_top") is not None:
        text_frame.margin_top = _pt(args["margin_top"])
    if args.get("margin_bottom") is not None:
        text_frame.margin_bottom = _pt(args["margin_bottom"])
    if args.get("word_wrap") is not None:
        text_frame.word_wrap = bool(args["word_wrap"])
    auto = args.get("text_autofit")
    if auto == "shrink_text":
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif auto == "grow_shape":
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    elif auto == "none":
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    valign = {"top": MSO_VERTICAL_ANCHOR.TOP, "middle": MSO_VERTICAL_ANCHOR.MIDDLE, "bottom": MSO_VERTICAL_ANCHOR.BOTTOM}
    if args.get("vertical_alignment") in valign:
        text_frame.vertical_anchor = valign[args["vertical_alignment"]]
    alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    font_color = _rgb(args.get("font_color"))
    for paragraph in text_frame.paragraphs:
        if args.get("alignment") in alignment:
            paragraph.alignment = alignment[args["alignment"]]
        if not paragraph.runs:
            paragraph.add_run()
        for run in paragraph.runs:
            font = run.font
            if args.get("font_name") is not None:
                font.name = str(args["font_name"])
            if args.get("font_size") is not None:
                font.size = _pt(args["font_size"])
            if args.get("bold") is not None:
                font.bold = bool(args["bold"])
            if args.get("italic") is not None:
                font.italic = bool(args["italic"])
            if font_color is not None:
                font.color.rgb = font_color


def _set_text(shape, text: object, args: dict) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.text = "" if text is None else str(text)
    _apply_text_frame(shape.text_frame, args)


SHAPES = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "round_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "ellipse": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "hexagon": MSO_SHAPE.HEXAGON,
    "pentagon": MSO_SHAPE.PENTAGON,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "chevron": MSO_SHAPE.CHEVRON,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "cloud": MSO_SHAPE.CLOUD,
    "arc": MSO_SHAPE.ARC,
    "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
    "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
}

CHARTS = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


def _shape_result(shape, slide_index: int) -> dict:
    return {
        "slide_index": slide_index,
        "shape_id": shape.shape_id,
        "shape_name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": shape.left.pt,
        "top": shape.top.pt,
        "width": shape.width.pt,
        "height": shape.height.pt,
    }


def action_status(args: dict) -> dict:
    state = _state()
    host = _presentation_host_info(args, state)
    path = Path(state["path"]) if state.get("path") else None
    return {
        **host,
        "connected_to_active_application": host["running_processes"] > 0,
        "active_presentation": bool(path and path.exists()),
        "managed_path": str(path) if path else None,
        "source_path": state.get("source_path"),
        "native_editable_output": True,
        "focus_policy": _focus_policy(),
        "live_semantics": "safe file-backed reload; edits are written as native PPTX objects then the managed deck is reopened in the selected presentation application",
    }


def action_capabilities(args: dict) -> dict:
    host = action_status(args)
    supported = {
        "text_box": "Shapes.add_textbox",
        "auto_shape": "Shapes.add_shape",
        "free_line_or_arrow": "Shapes.add_connector",
        "attached_connector": "file-backed native connector",
        "table": "Shapes.add_table",
        "chart": "Shapes.add_chart",
        "picture_or_svg": "Shapes.add_picture",
        "duplicate": "OOXML duplicate",
        "group": "Shapes.add_group_shape",
        "ungroup": "OOXML ungroup",
        "z_order": "OOXML tree order",
        "align": "exact geometry",
        "distribute": "exact geometry",
        "figure_audit": "deterministic OOXML geometry audit",
    }
    return {
        "detection": {
            "read_only": True,
            "launched_powerpoint": False,
            "active_deck_modified": False,
            "basis": ["python-pptx OOXML API", "Microsoft PowerPoint for Mac presence and AppleScript application metadata"],
        },
        "host": host,
        "native_object_families": [
            {"family": key, "powerpoint_api": value, "host_supported": True, "editable": key != "figure_audit"}
            for key, value in supported.items()
        ],
        "auto_shapes": [{"plugin_name": name, "value": int(value)} for name, value in SHAPES.items()] if args.get("include_auto_shapes", True) else [],
        "chart_types": [{"plugin_name": name, "value": int(value)} for name, value in CHARTS.items()] if args.get("include_chart_types", True) else [],
        "connector_types": [{"plugin_name": name} for name in ("straight", "elbow", "curve")],
        "arrowhead_styles": [{"plugin_name": name} for name in ("none", "open", "triangle", "stealth", "diamond", "oval")],
        "limitations": [
            "PowerPoint for Mac and WPS Presentation use an isolated file-backed working copy because they do not expose the Windows PowerPoint COM automation server used by this plugin.",
            "The managed deck is reopened after mutations; the default preserve focus policy uses background opening so the presentation does not repeatedly take over the desktop.",
            "Renderer exports use local LibreOffice/Poppler when available and remain separate from the editable PPTX.",
        ],
    }


def action_launch(args: dict) -> dict:
    state = _state()
    host_name, _ = _select_host(args, state)
    state["host_application"] = host_name
    supplied = args.get("file_path")
    if supplied:
        source = Path(supplied).expanduser().resolve()
        if not source.exists():
            raise FileNotFoundError(source)
        if args.get("read_only", False):
            working = source
            owned = False
        else:
            working = _managed_path(source.stem + "-working")
            shutil.copy2(source, working)
            owned = True
        state.update({"path": str(working), "source_path": str(source), "owned": owned, "metadata": {}})
    elif state.get("path") and Path(state["path"]).exists():
        working = Path(state["path"])
    elif args.get("create_if_missing", True):
        working = _managed_path()
        _new_presentation(working)
        state.update({"path": str(working), "source_path": None, "owned": True, "metadata": {}})
    else:
        raise RuntimeError("No managed presentation exists and create_if_missing=false.")
    _write_state(state)
    if args.get("visible", True):
        _refresh_presentation(working, state)
    return {**action_status({}), "opened_path": str(working), "working_copy": bool(state.get("owned"))}


def action_new_presentation(args: dict) -> dict:
    path = _managed_path()
    prs = _new_presentation(path)
    host_name, _ = _select_host(args)
    state = {"path": str(path), "source_path": None, "owned": True, "metadata": {}, "host_application": host_name}
    _write_state(state)
    _refresh_presentation(path, state)
    return {"path": str(path), "slide_count": len(prs.slides), "platform": sys.platform, "backend": "python-pptx-ooxml+application-reload", "host_application": host_name}


def _inventory(prs: Presentation, args: dict) -> list[dict]:
    result = []
    max_slides = int(args.get("max_slides", 100))
    max_shapes = int(args.get("max_shapes_per_slide", 200))
    for index, slide in enumerate(list(prs.slides)[:max_slides], 1):
        shapes = []
        for shape in list(slide.shapes)[:max_shapes]:
            item = _shape_result(shape, index)
            if args.get("include_text", True) and getattr(shape, "has_text_frame", False):
                item["text"] = shape.text
            item["is_picture"] = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            item["is_table"] = bool(getattr(shape, "has_table", False))
            item["is_chart"] = bool(getattr(shape, "has_chart", False))
            shapes.append(item)
        result.append({"slide_index": index, "shape_count": len(slide.shapes), "shapes": shapes})
    return result


def action_inspect(args: dict) -> dict:
    state, path, prs = _load()
    return {
        "path": str(path),
        "platform": sys.platform,
        "backend": "python-pptx-ooxml+application-reload",
        "host_application": state.get("host_application", "auto"),
        "slide_width": prs.slide_width.pt,
        "slide_height": prs.slide_height.pt,
        "slide_count": len(prs.slides),
        "slides": _inventory(prs, args),
    }


def action_add_slide(args: dict) -> dict:
    state, path, prs = _load()
    layout_index = {"blank": 6, "title": 0, "text": 1}.get(args.get("layout", "blank"), 6)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    position = args.get("position")
    if position and int(position) <= len(prs.slides):
        slide_ids = prs.slides._sldIdLst
        slide_id = slide_ids[-1]
        slide_ids.remove(slide_id)
        slide_ids.insert(max(0, int(position) - 1), slide_id)
    index = list(prs.slides).index(slide) + 1
    if args.get("name"):
        state.setdefault("metadata", {}).setdefault("slides", {})[str(index)] = {"name": args["name"]}
    _save(prs, state, path)
    return {"slide_index": index, "slide_count": len(prs.slides), "name": args.get("name")}


def action_activate_slide(args: dict) -> dict:
    state, path = _require_path(False)
    index = int(args["slide_index"])
    prs = Presentation(path)
    _slide(prs, index)
    _refresh_presentation(path, state, focus_policy="foreground")
    return {"slide_index": index, "activated": True, "note": "Explicit activation brought the managed deck forward; ordinary drawing preserves the user's current foreground application."}


def action_add_textbox(args: dict) -> dict:
    state, path, prs = _load()
    slide = _slide(prs, int(args["slide_index"]))
    shape = slide.shapes.add_textbox(_pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    _set_text(shape, args.get("text", ""), args)
    _set_fill(shape, args.get("fill_color"))
    _set_line(shape, args)
    _save(prs, state, path)
    return _shape_result(shape, int(args["slide_index"]))


def _shape_enum(args: dict):
    if args.get("shape_type_id") is not None:
        return MSO_SHAPE(int(args["shape_type_id"]))
    key = str(args.get("shape", "rectangle")).replace("msoShape", "").replace("-", "_").lower()
    return SHAPES.get(key, SHAPES.get(key.replace(" ", "_"), MSO_SHAPE.RECTANGLE))


def action_add_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    shape = slide.shapes.add_shape(_shape_enum(args), _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    if args.get("rotation") is not None:
        shape.rotation = float(args["rotation"])
    _set_fill(shape, args.get("fill_color"))
    _set_line(shape, args)
    if args.get("text") is not None:
        _set_text(shape, args["text"], args)
    _save(prs, state, path)
    return _shape_result(shape, index)


def action_add_image(args: dict) -> dict:
    state, path, prs = _load()
    if args.get("atomic_raster_unit") is not True or args.get("contains_reconstructable_content") is not False:
        raise ValueError("Mac PowerPoint images must be atomic_raster_unit=true and contains_reconstructable_content=false.")
    image_path = Path(args["image_path"]).expanduser().resolve()
    if not image_path.exists():
        raise FileNotFoundError(image_path)
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    shape = slide.shapes.add_picture(str(image_path), _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    for side in ("left", "top", "right", "bottom"):
        percent = args.get(f"crop_{side}_percent")
        if percent is not None:
            setattr(shape, f"crop_{side}", float(percent) / 100.0)
    state.setdefault("metadata", {}).setdefault("rasters", {})[str(shape.shape_id)] = {
        "slide_index": index,
        "shape_name": shape.name,
        "raster_reason": args["raster_reason"],
        "atomic_raster_unit": True,
        "contains_reconstructable_content": False,
        "decomposition_note": args["decomposition_note"],
        "source_is_tightly_cropped": bool(args.get("source_is_tightly_cropped")),
    }
    _save(prs, state, path)
    return {**_shape_result(shape, index), "raster_declaration": state["metadata"]["rasters"][str(shape.shape_id)]}


def _trim_line(x1: float, y1: float, x2: float, y2: float, start: float, end: float) -> tuple[float, float, float, float]:
    dx, dy = x2 - x1, y2 - y1
    length = math.hypot(dx, dy)
    if length <= start + end or length == 0:
        raise ValueError("Line clearances consume the entire segment.")
    ux, uy = dx / length, dy / length
    return x1 + ux * start, y1 + uy * start, x2 - ux * end, y2 - uy * end


def action_add_line(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    x1, y1, x2, y2 = _trim_line(
        float(args["begin_x"]), float(args["begin_y"]), float(args["end_x"]), float(args["end_y"]),
        float(args.get("start_clearance", 0)), float(args.get("end_clearance", 0)),
    )
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _pt(x1), _pt(y1), _pt(x2), _pt(y2))
    if args.get("name"):
        shape.name = args["name"]
    _set_line(shape, args)
    _set_arrow(shape, args.get("start_arrow", "none"), args.get("end_arrow", "none"))
    _save(prs, state, path)
    return _shape_result(shape, index)


def _site_point(shape, site: int) -> tuple[float, float]:
    left, top, width, height = shape.left.pt, shape.top.pt, shape.width.pt, shape.height.pt
    points = {
        1: (left + width / 2, top),
        2: (left + width, top + height / 2),
        3: (left + width / 2, top + height),
        4: (left, top + height / 2),
    }
    return points.get(((int(site) - 1) % 4) + 1)


def action_add_connector(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    source = _shape(slide, {"shape_name": args["source_name"]})
    target = _shape(slide, {"shape_name": args["target_name"]})
    x1, y1 = _site_point(source, int(args.get("source_site", 2)))
    x2, y2 = _site_point(target, int(args.get("target_site", 4)))
    connector_type = {"straight": MSO_CONNECTOR.STRAIGHT, "elbow": MSO_CONNECTOR.ELBOW, "curve": MSO_CONNECTOR.CURVE}[args.get("connector_type", "elbow")]
    shape = slide.shapes.add_connector(connector_type, _pt(x1), _pt(y1), _pt(x2), _pt(y2))
    if args.get("name"):
        shape.name = args["name"]
    _set_line(shape, args)
    _set_arrow(shape, args.get("start_arrow", "none"), args.get("end_arrow", "triangle"))
    _save(prs, state, path)
    return {**_shape_result(shape, index), "source_name": source.name, "target_name": target.name, "attachment_mode": "geometry-backed"}


def _style_cell(cell, args: dict) -> None:
    if args.get("text") is not None:
        cell.text = str(args["text"])
    _set_fill(cell, args.get("fill_color"))
    _apply_text_frame(cell.text_frame, args)
    margin = args.get("cell_margin")
    if margin is not None:
        cell.margin_left = cell.margin_right = cell.margin_top = cell.margin_bottom = _pt(margin)


def action_add_table(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    shape = slide.shapes.add_table(int(args["rows"]), int(args["columns"]), _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    data = args.get("data") or []
    table = shape.table
    for row in range(len(table.rows)):
        for col in range(len(table.columns)):
            cell_args = dict(args)
            if row < len(data) and col < len(data[row]):
                cell_args["text"] = "" if data[row][col] is None else data[row][col]
            if row < int(args.get("header_rows", 1)):
                cell_args["fill_color"] = args.get("header_fill_color", args.get("fill_color"))
                cell_args["font_color"] = args.get("header_font_color", args.get("font_color"))
                cell_args["bold"] = args.get("header_bold", True)
            elif args.get("banded_rows") and row % 2:
                cell_args["fill_color"] = args.get("band_fill_color", args.get("fill_color"))
            _style_cell(table.cell(row, col), cell_args)
    for override in args.get("cell_styles") or []:
        row, col = int(override["row"]) - 1, int(override["column"]) - 1
        _style_cell(table.cell(row, col), {**args, **override})
    _save(prs, state, path)
    return {**_shape_result(shape, index), "rows": len(table.rows), "columns": len(table.columns)}


def action_update_table_cell(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    if not getattr(shape, "has_table", False):
        raise ValueError(f"{shape.name} is not a table")
    row, col = int(args["row"]) - 1, int(args["column"]) - 1
    _style_cell(shape.table.cell(row, col), args)
    _save(prs, state, path)
    return {**_shape_result(shape, index), "row": row + 1, "column": col + 1, "text": shape.table.cell(row, col).text}


def action_update_table_layout(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    if not getattr(shape, "has_table", False):
        raise ValueError(f"{shape.name} is not a table")
    for column, width in zip(shape.table.columns, args.get("column_widths") or []):
        column.width = _pt(width)
    for row, height in zip(shape.table.rows, args.get("row_heights") or []):
        row.height = _pt(height)
    _save(prs, state, path)
    return {**_shape_result(shape, index), "column_widths": args.get("column_widths"), "row_heights": args.get("row_heights")}


def action_add_chart(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    data = CategoryChartData()
    data.categories = list(args["categories"])
    for series in args["series"]:
        data.add_series(str(series["name"]), tuple(series["values"]))
    chart_type = CHARTS.get(str(args.get("chart_type", "column_clustered")).replace("xl", "").lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    shape = slide.shapes.add_chart(chart_type, _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]), data)
    if args.get("name"):
        shape.name = args["name"]
    chart = shape.chart
    if args.get("title") is not None:
        chart.has_title = True
        chart.chart_title.text_frame.text = str(args["title"])
    chart.has_legend = bool(args.get("has_legend", True))
    if chart.has_legend:
        chart.legend.position = {
            "right": XL_LEGEND_POSITION.RIGHT,
            "left": XL_LEGEND_POSITION.LEFT,
            "top": XL_LEGEND_POSITION.TOP,
            "bottom": XL_LEGEND_POSITION.BOTTOM,
        }[args.get("legend_position", "right")]
    if args.get("chart_style") is not None:
        chart.chart_style = int(args["chart_style"])
    _save(prs, state, path)
    return {**_shape_result(shape, index), "chart_type": str(chart_type), "series_count": len(args["series"])}


def action_duplicate_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    source = _shape(slide, args)
    next_shape_id = slide.shapes._next_shape_id
    new_element = copy.deepcopy(source._element)
    c_nv_pr = new_element.xpath(".//p:cNvPr")[0]
    c_nv_pr.set("id", str(next_shape_id))
    c_nv_pr.set("name", args["new_name"])
    slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
    duplicate = list(slide.shapes)[-1]
    duplicate.name = args["new_name"]
    for prop in ("left", "top", "width", "height"):
        if args.get(prop) is not None:
            setattr(duplicate, prop, _pt(args[prop]))
    if args.get("rotation") is not None:
        duplicate.rotation = float(args["rotation"])
    _save(prs, state, path)
    return {**_shape_result(duplicate, index), "duplicated_from": source.name}


def action_group_shapes(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    members = [_shape(slide, {"shape_name": name}) for name in args["shape_names"]]
    group = slide.shapes.add_group_shape(members)
    if args.get("name"):
        group.name = args["name"]
    _save(prs, state, path)
    return {**_shape_result(group, index), "members": [shape.name for shape in group.shapes]}


def action_ungroup_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    group = _shape(slide, args)
    if shape_type_name(group) != "GROUP":
        raise ValueError(f"{group.name} is not a group")
    names = []
    for member in list(group.shapes):
        names.append(member.name)
        slide.shapes._spTree.insert_element_before(member._element, "p:extLst")
    group._element.getparent().remove(group._element)
    _save(prs, state, path)
    return {"slide_index": index, "ungrouped": args.get("shape_name") or args.get("shape_id"), "members": names}


def shape_type_name(shape) -> str:
    try:
        return shape.shape_type.name
    except Exception:
        return str(shape.shape_type)


def action_set_z_order(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    shape = _shape(slide, args)
    parent = shape._element.getparent()
    command = args["command"]
    repeat = int(args.get("repeat", 1))
    for _ in range(repeat):
        siblings = [item for item in parent if item.tag.endswith(("sp", "pic", "graphicFrame", "grpSp", "cxnSp"))]
        pos = siblings.index(shape._element)
        if command == "bring_to_front":
            parent.remove(shape._element)
            parent.insert_element_before(shape._element, "p:extLst")
        elif command == "send_to_back":
            parent.remove(shape._element)
            parent.insert(2, shape._element)
        elif command == "bring_forward" and pos + 1 < len(siblings):
            sibling = siblings[pos + 1]
            parent.remove(shape._element)
            parent.insert(parent.index(sibling) + 1, shape._element)
        elif command == "send_backward" and pos > 0:
            sibling = siblings[pos - 1]
            parent.remove(shape._element)
            parent.insert(parent.index(sibling), shape._element)
    _save(prs, state, path)
    return {**_shape_result(shape, index), "command": command}


def action_align_shapes(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    prs_slide = _slide(prs, index)
    shapes = [_shape(prs_slide, {"shape_name": name}) for name in args["shape_names"]]
    relative_slide = args.get("relative_to") == "slide"
    alignment = args["alignment"]
    if alignment == "left":
        value = 0 if relative_slide else min(s.left for s in shapes)
        for s in shapes: s.left = value
    elif alignment == "right":
        value = prs.slide_width if relative_slide else max(s.left + s.width for s in shapes)
        for s in shapes: s.left = value - s.width
    elif alignment == "center":
        value = prs.slide_width / 2 if relative_slide else sum(s.left + s.width / 2 for s in shapes) / len(shapes)
        for s in shapes: s.left = int(value - s.width / 2)
    elif alignment == "top":
        value = 0 if relative_slide else min(s.top for s in shapes)
        for s in shapes: s.top = value
    elif alignment == "bottom":
        value = prs.slide_height if relative_slide else max(s.top + s.height for s in shapes)
        for s in shapes: s.top = value - s.height
    elif alignment == "middle":
        value = prs.slide_height / 2 if relative_slide else sum(s.top + s.height / 2 for s in shapes) / len(shapes)
        for s in shapes: s.top = int(value - s.height / 2)
    _save(prs, state, path)
    return {"slide_index": index, "alignment": alignment, "shape_names": [s.name for s in shapes]}


def action_distribute_shapes(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    shapes = [_shape(slide, {"shape_name": name}) for name in args["shape_names"]]
    direction = args["direction"]
    if direction == "horizontal":
        shapes.sort(key=lambda s: s.left)
        start = 0 if args.get("relative_to") == "slide" else shapes[0].left
        end = prs.slide_width if args.get("relative_to") == "slide" else shapes[-1].left + shapes[-1].width
        gap = (end - start - sum(s.width for s in shapes)) / (len(shapes) - 1)
        cursor = start
        for shape in shapes:
            shape.left = int(cursor)
            cursor += shape.width + gap
    else:
        shapes.sort(key=lambda s: s.top)
        start = 0 if args.get("relative_to") == "slide" else shapes[0].top
        end = prs.slide_height if args.get("relative_to") == "slide" else shapes[-1].top + shapes[-1].height
        gap = (end - start - sum(s.height for s in shapes)) / (len(shapes) - 1)
        cursor = start
        for shape in shapes:
            shape.top = int(cursor)
            cursor += shape.height + gap
    _save(prs, state, path)
    return {"slide_index": index, "direction": direction, "shape_names": [s.name for s in shapes]}


def action_update_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    if args.get("new_name"):
        shape.name = args["new_name"]
    for prop in ("left", "top", "width", "height"):
        if args.get(prop) is not None:
            setattr(shape, prop, _pt(args[prop]))
    if args.get("rotation") is not None:
        shape.rotation = float(args["rotation"])
    if args.get("text") is not None:
        _set_text(shape, args["text"], args)
    elif getattr(shape, "has_text_frame", False):
        _apply_text_frame(shape.text_frame, args)
    _set_fill(shape, args.get("fill_color"))
    _set_line(shape, args)
    _save(prs, state, path)
    return _shape_result(shape, index)


def action_delete_shape(args: dict) -> dict:
    if args.get("confirm") is not True:
        raise ValueError("confirm=true is required")
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    deleted = {"shape_name": shape.name, "shape_id": shape.shape_id}
    shape._element.getparent().remove(shape._element)
    _save(prs, state, path)
    return {"slide_index": index, "deleted": deleted}


def action_audit_figure(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    findings = []
    max_findings = int(args.get("max_findings", 300))
    slide_area = max(1, prs.slide_width * prs.slide_height)
    raster_meta = state.get("metadata", {}).get("rasters", {})
    for shape in slide.shapes:
        if (shape.width <= 0 or shape.height <= 0) and shape.shape_type != MSO_SHAPE_TYPE.LINE:
            findings.append({"severity": "hard", "category": "geometry", "shape_name": shape.name, "message": "Shape has non-positive dimensions."})
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
            findings.append({"severity": "warning", "category": "bounds", "shape_name": shape.name, "message": "Shape extends beyond slide bounds."})
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ratio = (shape.width * shape.height) / slide_area
            if ratio >= float(args.get("large_raster_area_ratio", 0.08)) and str(shape.shape_id) not in raster_meta:
                findings.append({"severity": "hard", "category": "raster_editability", "shape_name": shape.name, "message": "Large picture has no atomic-raster declaration."})
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            font_size = 18.0
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        font_size = run.font.size.pt
                        break
            approx_capacity = max(1, (shape.width.pt / max(font_size * 0.55, 1)) * (shape.height.pt / max(font_size * 1.25, 1)))
            if len(shape.text) > approx_capacity * 1.3:
                findings.append({"severity": "warning", "category": "text_fit", "shape_name": shape.name, "message": "Text may overflow; verify the rendered slide."})
        if len(findings) >= max_findings:
            break
    hard_count = sum(1 for item in findings if item["severity"] == "hard")
    return {
        "slide_index": index,
        "path": str(path),
        "backend": "python-pptx-ooxml+application-reload",
        "host_application": state.get("host_application", "auto"),
        "shape_count": len(slide.shapes),
        "findings": findings,
        "hard_failure_count": hard_count,
        "warning_count": len(findings) - hard_count,
        "passed_deterministic_gate": hard_count == 0,
        "renderer_review_required": True,
    }


def _find_binary(name: str) -> str:
    binary = shutil.which(name)
    if not binary:
        raise RuntimeError(f"Required local renderer binary not found: {name}")
    return binary


def _render_pdf(source: Path, destination: Path) -> None:
    destination.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="scientific-illustrator-pdf-") as tmp:
        tmp_path = Path(tmp)
        _run([_find_binary("soffice"), "--headless", "--convert-to", "pdf", "--outdir", str(tmp_path), str(source)], timeout=180)
        generated = tmp_path / f"{source.stem}.pdf"
        if not generated.exists():
            candidates = list(tmp_path.glob("*.pdf"))
            if not candidates:
                raise RuntimeError("LibreOffice did not produce a PDF.")
            generated = candidates[0]
        shutil.copy2(generated, destination)


def action_export_slide_image(args: dict) -> dict:
    state, path = _require_path(False)
    output = Path(args["output_path"]).expanduser().resolve()
    if output.exists() and not args.get("overwrite", False):
        raise FileExistsError(f"Output exists: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    width, height, slide_index = int(args.get("width", 1920)), int(args.get("height", 1080)), int(args["slide_index"])
    with tempfile.TemporaryDirectory(prefix="scientific-illustrator-render-") as tmp:
        pdf_path = Path(tmp) / "deck.pdf"
        _render_pdf(path, pdf_path)
        jpeg = output.suffix.lower() in (".jpg", ".jpeg")
        rendered = Path(tmp) / ("slide.jpg" if jpeg else "slide.png")
        format_flag = "-jpeg" if jpeg else "-png"
        _run([_find_binary("pdftoppm"), "-f", str(slide_index), "-l", str(slide_index), "-singlefile", format_flag, "-scale-to-x", str(width), "-scale-to-y", str(height), str(pdf_path), str(rendered.with_suffix(""))], timeout=180)
        shutil.copy2(rendered, output)
        mime = "image/jpeg" if jpeg else "image/png"
    return {"slide_index": slide_index, "output_path": str(output), "width": width, "height": height, "mime_type": mime, "renderer": "LibreOffice+Poppler"}


def action_save(args: dict) -> dict:
    state, path = _require_path(False)
    output_raw = args.get("output_path")
    if not output_raw:
        return {"output_path": str(path), "format": "pptx", "saved": True}
    output = Path(output_raw).expanduser().resolve()
    if output.exists() and not args.get("overwrite", False):
        raise FileExistsError(f"Output exists: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    fmt = args.get("format") or ("pdf" if output.suffix.lower() == ".pdf" else "pptx")
    if fmt == "pdf":
        _render_pdf(path, output)
    else:
        shutil.copy2(path, output)
        state["last_saved_output"] = str(output)
        _write_state(state)
    return {"output_path": str(output), "format": fmt, "saved": True}


def action_close_presentation(args: dict) -> dict:
    if args.get("confirm") is not True:
        raise ValueError("confirm=true is required")
    state, path = _require_path(False)
    if args.get("save_changes") == "save" and args.get("output_path"):
        action_save({"output_path": args["output_path"], "overwrite": args.get("overwrite", False), "format": "pptx"})
    host_name, _ = _select_host(state=state)
    if sys.platform != "darwin" or host_name != "powerpoint":
        return {
            "closed": False,
            "detached": True,
            "path": str(path),
            "file_retained": path.exists(),
            "note": "The file-backed WPS backend does not force-close application windows; close the managed deck in WPS when convenient.",
        }
    script = r'''
on run argv
    set targetName to item 1 of argv
    tell application "Microsoft PowerPoint"
        repeat with deck in presentations
            try
                if (name of deck as text) is targetName then close deck saving no
            end try
        end repeat
    end tell
end run
'''
    _osascript(script, path.name, check=False)
    return {"closed": True, "path": str(path), "file_retained": path.exists()}


def action_quit_application(args: dict) -> dict:
    if args.get("confirm") is not True:
        raise ValueError("confirm=true is required")
    state = _state()
    host = _presentation_host_info(args, state)
    if host["host_application"] != "powerpoint" or sys.platform != "darwin":
        raise RuntimeError("Safe programmatic application quit is only available for Microsoft PowerPoint on macOS in the file-backed backend.")
    if host["presentation_count"] != 0:
        raise RuntimeError("PowerPoint still has open presentations; close them explicitly before quitting.")
    expected = int(args.get("expected_process_id", 0))
    if expected not in host["process_ids"]:
        raise RuntimeError("expected_process_id does not match the current PowerPoint process.")
    _osascript('tell application "Microsoft PowerPoint" to quit', check=True)
    return {"quit": True, "process_id": expected}


ACTIONS = {
    "status": action_status,
    "capabilities": action_capabilities,
    "launch": action_launch,
    "new_presentation": action_new_presentation,
    "inspect": action_inspect,
    "audit_figure": action_audit_figure,
    "activate_slide": action_activate_slide,
    "add_slide": action_add_slide,
    "add_textbox": action_add_textbox,
    "add_shape": action_add_shape,
    "add_image": action_add_image,
    "add_line": action_add_line,
    "add_connector": action_add_connector,
    "add_table": action_add_table,
    "update_table_cell": action_update_table_cell,
    "update_table_layout": action_update_table_layout,
    "add_chart": action_add_chart,
    "duplicate_shape": action_duplicate_shape,
    "group_shapes": action_group_shapes,
    "ungroup_shape": action_ungroup_shape,
    "set_z_order": action_set_z_order,
    "align_shapes": action_align_shapes,
    "distribute_shapes": action_distribute_shapes,
    "update_shape": action_update_shape,
    "delete_shape": action_delete_shape,
    "export_slide_image": action_export_slide_image,
    "save": action_save,
    "close_presentation": action_close_presentation,
    "quit_application": action_quit_application,
}


def main() -> None:
    if len(sys.argv) != 2:
        raise SystemExit("Usage: powerpoint-mac-bridge.py <payload-base64>")
    payload = json.loads(base64.b64decode(sys.argv[1]).decode("utf-8"))
    action = str(payload["action"])
    arguments = payload.get("arguments") or {}
    if action not in ACTIONS:
        raise ValueError(f"Unsupported macOS PowerPoint action: {action}")
    result = ACTIONS[action](arguments)
    sys.stdout.write(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        sys.stderr.write(f"{type(exc).__name__}: {exc}\n")
        raise
